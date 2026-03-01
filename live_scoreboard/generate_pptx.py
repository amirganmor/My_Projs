"""Generate a 12-slide professional PowerPoint for the Scoreboard platform."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARKER_BG = RGBColor(0x12, 0x12, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0xA5, 0x00)
MUTED = RGBColor(0x99, 0x99, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCREENSHOTS_DIR = os.path.join(os.path.dirname(__file__), "screenshots")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_rect(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_screenshot(slide, filename, left, top, width, height=None):
    path = os.path.join(SCREENSHOTS_DIR, filename)
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, left, top, width, height)
        else:
            slide.shapes.add_picture(path, left, top, width=width)


def accent_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


# ── SLIDE 1: Title ──────────────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARKER_BG)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "Real-Time Multi-Sport Scoreboard", font_size=44, bold=True, color=WHITE,
         alignment=PP_ALIGN.CENTER)

tf = add_text(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
              "A Streaming Data Engineering Platform", font_size=24, color=ACCENT_BLUE,
              alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "NBA  |  NFL  |  Premier League  |  La Liga  |  Champions League",
         font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.5),
              "", font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)
tf.paragraphs[0].text = "Apache Kafka  +  PySpark Structured Streaming  +  PostgreSQL  +  Flask  +  Docker Compose"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = MUTED
add_para(tf, "14 services  |  5 Kafka topics  |  4 database tables  |  3 parallel pipelines",
         font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

accent_line(slide, Inches(3.6))

# ── SLIDE 2: Problem & Goals ────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Problem & Goals", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

problems = [
    "Sports data generated in real time across multiple leagues, APIs, and formats",
    "Different consumers need different views: fans want live scores, analysts want SQL-queryable data",
    "Must decouple producers from consumers and handle bursty concurrent updates",
    "APIs have rate limits -- need caching, deduplication, and smart polling",
]
tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3),
              "The Problem", font_size=20, bold=True, color=WHITE)
for p in problems:
    add_para(tf, f"  {p}", font_size=14, color=LIGHT_GRAY)

add_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5))
tf = add_text(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(4.5),
              "Three Parallel Pipelines", font_size=20, bold=True, color=ACCENT_GREEN)
pipelines = [
    ("Pipeline 1 -- Live Scoreboard UI", "Sources -> Kafka -> Flask -> Browser (3s refresh)"),
    ("Pipeline 2 -- Streaming Analytics", "Sources -> Kafka -> PySpark -> PostgreSQL (3 tables)"),
    ("Pipeline 3 -- Batch Player Stats", "ESPN/nba_api -> psycopg2 -> PostgreSQL (every 6h)"),
]
for title, desc in pipelines:
    add_para(tf, f"  {title}", font_size=15, bold=True, color=WHITE, space_before=Pt(14))
    add_para(tf, f"    {desc}", font_size=13, color=LIGHT_GRAY, space_before=Pt(2))

# ── SLIDE 3: Leagues & Data Sources ─────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Supported Leagues & Data Sources", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

leagues = [
    ("NBA", "Basketball", "nba_api Python SDK + ESPN", "Free"),
    ("NFL", "American Football", "ESPN REST API", "Free"),
    ("Premier League", "Soccer (England)", "ESPN REST API", "Free"),
    ("La Liga", "Soccer (Spain)", "ESPN REST API", "Free"),
    ("Champions League", "Soccer (Europe)", "ESPN REST API", "Free"),
]

headers = ["League", "Sport", "Data Source", "Cost"]
col_lefts = [Inches(1.0), Inches(3.3), Inches(5.8), Inches(10.0)]
col_widths = [Inches(2.3), Inches(2.5), Inches(4.2), Inches(2.0)]

for i, h in enumerate(headers):
    add_rect(slide, col_lefts[i], Inches(1.6), col_widths[i], Inches(0.45))
    add_text(slide, col_lefts[i] + Inches(0.1), Inches(1.63), col_widths[i], Inches(0.4),
             h, font_size=14, bold=True, color=ACCENT_BLUE)

for row_idx, (league, sport, source, cost) in enumerate(leagues):
    y = Inches(2.15) + Inches(row_idx * 0.55)
    vals = [league, sport, source, cost]
    for i, val in enumerate(vals):
        c = ACCENT_GREEN if i == 3 else LIGHT_GRAY
        add_text(slide, col_lefts[i] + Inches(0.1), y, col_widths[i], Inches(0.4),
                 val, font_size=13, color=c)

add_rect(slide, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.8))
tf = add_text(slide, Inches(1.3), Inches(5.3), Inches(10.8), Inches(1.6),
              "All data sources are FREE with no API keys required", font_size=16, bold=True, color=ACCENT_GREEN)
add_para(tf, "ESPN endpoints: site.api.espn.com/apis/site/v2/sports/{sport}/{league}/scoreboard", font_size=12, color=MUTED)
add_para(tf, "Match details: .../summary?event={id} -- goals, cards, subs, commentary, lineups, 20+ stats", font_size=12, color=MUTED)
add_para(tf, "Dynamic cache TTL: 3s for live soccer, 5s for live NFL, 30s for finished games", font_size=12, color=MUTED)

# ── SLIDE 4: Architecture ───────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "System Architecture", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

arch_text = """                        DATA SOURCES
    nba_api       ESPN NFL      ESPN PL      ESPN Liga     ESPN UCL
       |              |              |              |              |
       v              v              v              v              v
                   KAFKA PRODUCERS (Python, 3-10s intervals)
                   Enriched messages: scores + events + stats + lineups
       |              |              |              |              |
       v              v              v              v              v
  [nba_scores]  [nfl_scores]  [pl_scores]  [liga_scores]  [cl_scores]
                              |
              +---------------+----------------+
              v                                v
   FLASK SCOREBOARD UI              PYSPARK STRUCTURED STREAMING
   + Match Trackers                 foreachBatch -> 3 JDBC tables
   + Tabs: Live/Stats/Lineups            |
              |                           v
              v                     PostgreSQL :5433
   Web Browser :5050                (game_scores, match_events, period_scores)

              PLAYER STATS BATCH FETCHER (every 6 hours)
              nba_api + ESPN -> psycopg2 -> PostgreSQL (player_season_stats)"""

add_rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8))
add_text(slide, Inches(0.8), Inches(1.5), Inches(12), Inches(5.5),
         arch_text, font_size=11, color=ACCENT_GREEN, font_name="Courier New")

# ── SLIDE 5: Technology Stack ────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Technology Stack", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

tech = [
    ("Python 3.11", "All producers, UI, Spark job, stats fetcher"),
    ("Apache Kafka 3.7.0 (KRaft)", "Event streaming, topic-based pub/sub, no ZooKeeper"),
    ("Apache Spark 3.5.0", "Structured Streaming, micro-batch to PostgreSQL"),
    ("PostgreSQL 16", "Persistent analytics store (4 tables)"),
    ("Flask 3.x", "Scoreboard UI backend with background Kafka consumer"),
    ("kafka-python-ng 2.x", "Producer + consumer library with retry/backoff"),
    ("Docker Compose v2", "Full-stack orchestration -- 14 services"),
    ("Kafka UI 0.7.2", "Topic and consumer group monitoring"),
]

for i, (name, desc) in enumerate(tech):
    y = Inches(1.6) + Inches(i * 0.65)
    add_rect(slide, Inches(0.8), y, Inches(4), Inches(0.55))
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(3.8), Inches(0.45),
             name, font_size=15, bold=True, color=ACCENT_GREEN)
    add_text(slide, Inches(5.2), y + Inches(0.05), Inches(7.5), Inches(0.45),
             desc, font_size=14, color=LIGHT_GRAY)

# ── SLIDE 6: Kafka Layer & Message Schema ────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Kafka Layer & Enriched Message Schema", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

topics = ["nba_scores", "nfl_scores", "premier_league_scores", "la_liga_scores", "champions_league_scores"]
tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(2),
              "5 Kafka Topics", font_size=20, bold=True, color=WHITE)
for t in topics:
    add_para(tf, f"  {t}", font_size=14, color=ACCENT_GREEN)

add_para(tf, "", font_size=8)
add_para(tf, "Fingerprint Deduplication", font_size=18, bold=True, color=WHITE, space_before=Pt(16))
add_para(tf, "  Only publishes when scores, statuses,", font_size=13, color=LIGHT_GRAY)
add_para(tf, "  or play counts actually change", font_size=13, color=LIGHT_GRAY)
add_para(tf, "  Avoids flooding Kafka with identical data", font_size=13, color=LIGHT_GRAY)

schema = """{
  "league": "premier_league",
  "games": [{
    "gameId": "672341",
    "statusText": "Second Half - 67'",
    "homeTeam": {"teamName": "Liverpool", "score": 2},
    "awayTeam": {"teamName": "Man City", "score": 1},
    "periods": [...],
    "events": [{"type":"goal","minute":23,
                "player":"Salah","team":"Liverpool"}],
    "stats": {"possession":["62.3","37.7"],
              "shots":["15","7"], ...20+ fields},
    "plays": [{"minute":"67'",
               "text":"Goal! Salah header...",
               "type":"goal"}],
    "lineups": [{"team":"Liverpool",
                 "color":"#C8102E",
                 "formation":"4-3-3",
                 "players":[...]}]
  }]
}"""
add_rect(slide, Inches(5.8), Inches(1.3), Inches(6.8), Inches(5.8))
add_text(slide, Inches(6.0), Inches(1.4), Inches(6.5), Inches(5.5),
         schema, font_size=11, color=ACCENT_GREEN, font_name="Courier New")

# ── SLIDE 7: Scoreboard UI ──────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Live Scoreboard UI", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

add_screenshot(slide, "01_main_scoreboard.png", Inches(0.5), Inches(1.5), Inches(6.2))
add_screenshot(slide, "03_after_wait.png", Inches(6.9), Inches(1.5), Inches(6.0))

tf = add_text(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(1.2),
              "", font_size=14, color=LIGHT_GRAY)
tf.paragraphs[0].text = "Dark-themed web interface at localhost:5050  |  Auto-refreshes every 3 seconds  |  Click any match for live tracker"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY

# ── SLIDE 8: Live Match Tracker ──────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Live Match Tracker Deep Dive", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

add_screenshot(slide, "02_live_tracker.png", Inches(0.5), Inches(1.5), Inches(6))
add_screenshot(slide, "05_nba_live_tracker.png", Inches(6.8), Inches(1.5), Inches(6))

features = [
    "SVG pitch/court/field with animated ball (BallController)",
    "Smooth CSS cubic-bezier transitions between positions",
    "Possession indicator, momentum timeline, live commentary feed",
    "Text-based position inference from ESPN commentary data",
]
tf = add_text(slide, Inches(0.8), Inches(5.2), Inches(7), Inches(2),
              "", font_size=13, color=LIGHT_GRAY)
for f in features:
    if tf.paragraphs[0].text == "":
        tf.paragraphs[0].text = f"  {f}"
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    else:
        add_para(tf, f"  {f}", font_size=13, color=LIGHT_GRAY)

add_rect(slide, Inches(8), Inches(5.2), Inches(4.8), Inches(1.8))
tf = add_text(slide, Inches(8.2), Inches(5.3), Inches(4.5), Inches(1.6),
              "Free Data Limitation", font_size=15, bold=True, color=ACCENT_ORANGE)
add_para(tf, "ESPN text data: ~60-70% accuracy", font_size=12, color=LIGHT_GRAY)
add_para(tf, "Paid feeds (Sportradar/Opta):", font_size=12, color=LIGHT_GRAY)
add_para(tf, "  100% accuracy, sub-second updates", font_size=12, color=ACCENT_GREEN)

# ── SLIDE 9: Spark Streaming ────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Spark Streaming to PostgreSQL", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

tf = add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4),
              "Pipeline Stages", font_size=20, bold=True, color=WHITE)
stages = [
    ("1. Read", "Spark Structured Streaming subscribes to all 5 Kafka topics"),
    ("2. Parse", "JSON deserialized with enhanced schema (games + events + periods)"),
    ("3. Transform", "explode() flattens into rows for each destination table"),
    ("4. Write", "foreachBatch writes to 3 PostgreSQL tables via JDBC"),
]
for title, desc in stages:
    add_para(tf, f"  {title}", font_size=15, bold=True, color=ACCENT_GREEN, space_before=Pt(12))
    add_para(tf, f"    {desc}", font_size=13, color=LIGHT_GRAY, space_before=Pt(2))

tables_data = [
    ("game_scores", "Score snapshots (append-only time series)"),
    ("match_events", "Goals, cards, substitutions (soccer)"),
    ("period_scores", "Quarter/half score breakdowns"),
]
add_rect(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3))
tf = add_text(slide, Inches(7.3), Inches(1.6), Inches(5), Inches(2.5),
              "Target Tables", font_size=18, bold=True, color=ACCENT_GREEN)
for tbl, desc in tables_data:
    add_para(tf, f"  {tbl}", font_size=14, bold=True, color=WHITE, space_before=Pt(10))
    add_para(tf, f"    {desc}", font_size=12, color=LIGHT_GRAY, space_before=Pt(2))

config_items = [
    ("Trigger", "processingTime = 10 seconds"),
    ("Checkpoint", "/tmp/spark_checkpoint/scoreboard_v2"),
    ("Spark Master", "spark://spark-master:7077"),
    ("JARs", "Pre-downloaded at Docker build time (no Maven at runtime)"),
]
add_rect(slide, Inches(7), Inches(4.8), Inches(5.5), Inches(2.2))
tf = add_text(slide, Inches(7.3), Inches(4.9), Inches(5), Inches(2),
              "Configuration", font_size=16, bold=True, color=ACCENT_BLUE)
for k, v in config_items:
    add_para(tf, f"  {k}: {v}", font_size=12, color=LIGHT_GRAY)

# ── SLIDE 10: Code Quality ──────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Code Quality & Engineering Practices", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

quality_left = [
    ("Shared Soccer Engine (DRY)", "3 league producers are 12-line wrappers around shared/soccer_producer.py -- eliminates ~500 lines of duplication"),
    ("Structured Logging", "All Python files use the logging module with timestamps and severity levels instead of print()"),
    ("Docstrings & Constants", "Every public function documented; magic numbers extracted to named constants (MAX_DRIVES, SUMMARY_TTL_LIVE, etc.)"),
    ("Pinned Dependencies", "All 8 requirements.txt files use version range pins (>=x,<y) for reproducible builds"),
]

quality_right = [
    ("HTML Escaping (XSS)", "escapeHtml() applied to 40+ dynamic content insertions in the frontend"),
    ("Kafka Retry with Backoff", "Consumer wraps in while-True with exponential backoff (3s to 30s max)"),
    ("Specific Exception Handling", "requests.RequestException, json.JSONDecodeError, KeyError instead of bare except"),
    ("Docker Image Pinning", "kafka:3.7.0, kafka-ui:v0.7.2, postgres:16, spark:3.5.0 -- no :latest tags"),
]

for i, (title, desc) in enumerate(quality_left):
    y = Inches(1.5) + Inches(i * 1.15)
    add_rect(slide, Inches(0.5), y, Inches(5.8), Inches(1.0))
    tf = add_text(slide, Inches(0.7), y + Inches(0.05), Inches(5.5), Inches(0.9),
                  title, font_size=14, bold=True, color=ACCENT_GREEN)
    add_para(tf, desc, font_size=11, color=LIGHT_GRAY, space_before=Pt(2))

for i, (title, desc) in enumerate(quality_right):
    y = Inches(1.5) + Inches(i * 1.15)
    add_rect(slide, Inches(6.8), y, Inches(5.8), Inches(1.0))
    tf = add_text(slide, Inches(7.0), y + Inches(0.05), Inches(5.5), Inches(0.9),
                  title, font_size=14, bold=True, color=ACCENT_GREEN)
    add_para(tf, desc, font_size=11, color=LIGHT_GRAY, space_before=Pt(2))

add_text(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
         "Code Review Score: 8.5 / 10", font_size=22, bold=True, color=ACCENT_GREEN,
         alignment=PP_ALIGN.CENTER)

# ── SLIDE 11: Production & Future ────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "Production Scaling & Future Work", font_size=36, bold=True, color=ACCENT_BLUE)
accent_line(slide, Inches(1.1))

prod_items = [
    ("Paid Data Providers", "Sportradar, Opta, Genius Sports for real-time x,y coordinate tracking and sub-second updates"),
    ("WebSocket / SSE", "Replace HTTP polling with push-based updates for lower latency delivery to browser"),
    ("Redis State Store", "Cache live match state for instant reads, pub/sub fanout to multiple UI instances"),
    ("Grafana Dashboards", "Connect to PostgreSQL for real-time analytics visualization and alerting"),
]

future_items = [
    ("ML Predictions", "Win probability models trained on historical score progression data"),
    ("Horizontal Scaling", "Kafka partitions, Spark workers, load-balanced Flask instances"),
    ("CI/CD Pipeline", "GitHub Actions: lint, test, build images, push to container registry"),
    ("More Leagues", "Bundesliga, Serie A, MLS, MLB, NHL -- architecture supports any league with a public API"),
]

tf = add_text(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(0.5),
              "Production Requirements", font_size=18, bold=True, color=ACCENT_ORANGE)
for i, (title, desc) in enumerate(prod_items):
    y = Inches(2.0) + Inches(i * 1.1)
    add_rect(slide, Inches(0.5), y, Inches(6), Inches(0.95))
    tf2 = add_text(slide, Inches(0.7), y + Inches(0.05), Inches(5.7), Inches(0.85),
                   title, font_size=14, bold=True, color=ACCENT_GREEN)
    add_para(tf2, desc, font_size=11, color=LIGHT_GRAY, space_before=Pt(2))

tf = add_text(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(0.5),
              "Future Development", font_size=18, bold=True, color=ACCENT_BLUE)
for i, (title, desc) in enumerate(future_items):
    y = Inches(2.0) + Inches(i * 1.1)
    add_rect(slide, Inches(6.8), y, Inches(6), Inches(0.95))
    tf2 = add_text(slide, Inches(7.0), y + Inches(0.05), Inches(5.7), Inches(0.85),
                   title, font_size=14, bold=True, color=ACCENT_GREEN)
    add_para(tf2, desc, font_size=11, color=LIGHT_GRAY, space_before=Pt(2))

add_rect(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1))
tf = add_text(slide, Inches(0.8), Inches(6.15), Inches(12), Inches(0.9),
              "Key Insight: The #1 bottleneck for live tracker quality is the data source, not the code.",
              font_size=16, bold=True, color=ACCENT_ORANGE)
add_para(tf, "Free ESPN text data achieves ~60-70% tracker accuracy. Paid coordinate feeds (Sportradar/Opta) would bring it to 100%.",
         font_size=13, color=LIGHT_GRAY)

# ── SLIDE 12: Summary & Q&A ─────────────────────────────────────────────────

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARKER_BG)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Summary", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

takeaways = [
    "Event-driven architecture with Kafka decouples 5 producers from 3 consumers",
    "Unified enriched JSON schema normalizes heterogeneous data sources",
    "Live match tracker with animated ball, possession, momentum, and commentary",
    "PySpark Structured Streaming writes to PostgreSQL in 10-second micro-batches",
    "Code quality score: 8.5/10 -- logging, docstrings, DRY, pinned deps, XSS prevention",
    "14 Docker Compose services, single command to run the entire platform",
    "Free data sources reach ~60-70% tracker accuracy; paid feeds unlock 100%",
]

tf = add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(3.5),
              "", font_size=16, color=LIGHT_GRAY)
for i, t in enumerate(takeaways):
    if i == 0:
        tf.paragraphs[0].text = f"  {t}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    else:
        add_para(tf, f"  {t}", font_size=16, color=LIGHT_GRAY, space_before=Pt(10))

add_rect(slide, Inches(2), Inches(5.2), Inches(9), Inches(1.6))
tf = add_text(slide, Inches(2.3), Inches(5.3), Inches(8.5), Inches(1.4),
              "How to Run", font_size=18, bold=True, color=ACCENT_GREEN)
add_para(tf, "  docker compose up -d --build", font_size=15, color=ACCENT_BLUE, space_before=Pt(6))
add_para(tf, "  Scoreboard: localhost:5050  |  Kafka UI: localhost:8080  |  PostgreSQL: localhost:5433", font_size=13, color=LIGHT_GRAY)
add_para(tf, "  GitHub: github.com/dreamgroup-il/nba-proj", font_size=13, color=MUTED)

# ── SAVE ─────────────────────────────────────────────────────────────────────

output_path = os.path.join(os.path.dirname(__file__), "scoreboard_presentation.pptx")
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to {output_path}")
